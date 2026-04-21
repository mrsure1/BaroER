# -*- coding: utf-8 -*-
"""
바로응급실 · 프로젝트 소개 PPTX 생성기
=====================================

python-pptx 로 18장짜리 발표 덱을 만든다. 교실·강의실에서 그대로 띄울 수
있도록 **16:9 와이드스크린**, **맑은 고딕(Windows 기본 한글 폰트)** 을
쓰고 CDN 이나 외부 이미지 의존성은 전혀 없다.

디자인 토큰 (앱의 globals.css 와 동일)
  - Primary  : #E53935 (응급 레드)
  - Accent   : #0D9488 (medical teal)
  - Text     : #101218 / muted #636573 / subtle #8E91A0
  - Surface  : #FAFAFB / surface-2 #F4F4F7 / border #E7E7ED

실행
  python scripts/generate_presentation.py
결과
  docs/BaroER_발표자료.pptx  (원클릭 배포 가능)
"""

from __future__ import annotations

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------------------------------------------------------------------------
# Design tokens
# ---------------------------------------------------------------------------

PRIMARY       = RGBColor(0xE5, 0x39, 0x35)
PRIMARY_DARK  = RGBColor(0xD3, 0x2F, 0x2F)
PRIMARY_SOFT  = RGBColor(0xFD, 0xEB, 0xEB)
ACCENT        = RGBColor(0x0D, 0x94, 0x88)
ACCENT_SOFT   = RGBColor(0xDC, 0xFC, 0xF7)
TEXT          = RGBColor(0x10, 0x12, 0x18)
TEXT_MUTED    = RGBColor(0x63, 0x65, 0x73)
TEXT_SUBTLE   = RGBColor(0x8E, 0x91, 0xA0)
BORDER        = RGBColor(0xE7, 0xE7, 0xED)
BORDER_STRONG = RGBColor(0xD2, 0xD2, 0xDC)
SURFACE       = RGBColor(0xFA, 0xFA, 0xFB)
SURFACE2      = RGBColor(0xF4, 0xF4, 0xF7)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
STATUS_OK     = RGBColor(0x22, 0xC5, 0x5E)
STATUS_BUSY   = RGBColor(0xF9, 0x73, 0x22)
STATUS_FULL   = RGBColor(0xEF, 0x44, 0x44)
KTAS_1        = RGBColor(0x00, 0x67, 0xB3)   # 파랑 · 소생
KTAS_2        = PRIMARY                         # 빨강 · 긴급
KTAS_3        = RGBColor(0xF5, 0x9E, 0x0B)   # 주황 · 응급
KTAS_4        = RGBColor(0x22, 0xC5, 0x5E)   # 초록 · 준응급
KTAS_5        = RGBColor(0xA1, 0xA1, 0xAA)   # 회색 · 비응급

FONT_KR = "맑은 고딕"  # Windows 기본, 거의 모든 발표 환경에서 사용 가능
FONT_EN = "Segoe UI"

# 16:9 · 기본 13.333 × 7.5 inch
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75, corner=None):
    """rounded 또는 rectangle 셰이프 추가."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if corner else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    if corner is not None:
        # adjustments[0] 은 라운드 비율 (0~0.5)
        try:
            shape.adjustments[0] = corner
        except Exception:
            pass
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    return shape


def add_text(
    slide, x, y, w, h, text,
    *, size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP, font=FONT_KR, line_spacing=1.2,
    italic=False, letter_spacing=None,
):
    """하나의 텍스트 상자에 한 줄 또는 여러 줄 텍스트."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_rich_line(
    slide, x, y, w, h, runs,
    *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.25,
):
    """(텍스트, {size,bold,color,font}) 튜플 리스트로 다중 서식 한 줄."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for text, opts in runs:
        run = p.add_run()
        run.text = text
        run.font.name = opts.get("font", FONT_KR)
        run.font.size = Pt(opts.get("size", 14))
        run.font.bold = opts.get("bold", False)
        run.font.italic = opts.get("italic", False)
        run.font.color.rgb = opts.get("color", TEXT)
    return tb


def add_line(slide, x1, y1, x2, y2, color=BORDER_STRONG, weight=0.75):
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def add_pill(slide, x, y, w, h, text, *, fill=PRIMARY_SOFT, color=PRIMARY, size=10, bold=True, border=None):
    pill = add_rect(slide, x, y, w, h, fill=fill, line=border, corner=0.5)
    add_text(
        slide, x, y, w, h, text,
        size=size, bold=bold, color=color,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    return pill


def add_circle(slide, x, y, d, fill, line=None):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    c.fill.solid()
    c.fill.fore_color.rgb = fill
    if line is None:
        c.line.fill.background()
    else:
        c.line.color.rgb = line
        c.line.width = Pt(1)
    c.shadow.inherit = False
    return c


# ---------------------------------------------------------------------------
# Slide shell
# ---------------------------------------------------------------------------

def draw_shell(slide, *, index, total, kicker=None, accent_tone="primary"):
    """모든 슬라이드의 공통 프레임(좌측 액센트 바 + 상단 kicker + 페이지 카운터)."""
    # 좌측 세로 액센트 바
    if accent_tone != "none":
        tone = PRIMARY if accent_tone == "primary" else ACCENT
        add_rect(
            slide,
            Inches(0.45), Inches(1.55), Inches(0.08), Inches(1.4),
            fill=tone, corner=0.5,
        )

    if kicker:
        add_rich_line(
            slide, Inches(0.6), Inches(0.42), Inches(10), Inches(0.3),
            [
                ("—  ", {"size": 10, "bold": True, "color": TEXT_SUBTLE}),
                (kicker.upper(), {"size": 10, "bold": True, "color": TEXT_SUBTLE}),
            ],
        )

    # 페이지 카운터 (우상단)
    add_text(
        slide, Inches(12.0), Inches(0.42), Inches(1.0), Inches(0.3),
        f"{index:02d}  /  {total:02d}",
        size=10, bold=True, color=TEXT_SUBTLE, align=PP_ALIGN.RIGHT, font=FONT_EN,
    )

    # 푸터 — 브랜드 + 타이틀 작은 글자
    add_rect(
        slide, Inches(0), SLIDE_H - Inches(0.45), SLIDE_W, Inches(0.01),
        fill=BORDER,
    )
    add_text(
        slide, Inches(0.6), SLIDE_H - Inches(0.4), Inches(6), Inches(0.3),
        "바로응급실 · BaroER",
        size=9, bold=True, color=PRIMARY,
    )
    add_text(
        slide, Inches(6.8), SLIDE_H - Inches(0.4), Inches(6), Inches(0.3),
        "Emergency · Real-time · Mobile   —   Project Deck 2026",
        size=9, color=TEXT_SUBTLE, align=PP_ALIGN.RIGHT, font=FONT_EN,
    )


def draw_title(slide, title, highlight=None, eyebrow=None, y=Inches(1.55)):
    """타이틀 + 강조 그라데이션(단색 fallback) + eyebrow 뱃지."""
    top = y
    if eyebrow:
        add_pill(
            slide, Inches(0.8), top, Inches(2.6), Inches(0.38),
            eyebrow, fill=PRIMARY_SOFT, color=PRIMARY, size=11, bold=True,
        )
        top += Inches(0.5)

    # 타이틀
    add_text(
        slide, Inches(0.8), top, Inches(11.7), Inches(0.9), title,
        size=36, bold=True, color=TEXT, line_spacing=1.08,
    )
    top += Inches(0.85)

    if highlight:
        add_text(
            slide, Inches(0.8), top, Inches(11.7), Inches(0.9), highlight,
            size=36, bold=True, color=PRIMARY, line_spacing=1.08,
        )
        top += Inches(0.85)
    return top


def draw_subtitle(slide, text, y, max_width=Inches(11.0)):
    add_text(
        slide, Inches(0.8), y, max_width, Inches(1.0), text,
        size=15, color=TEXT_MUTED, line_spacing=1.45,
    )
    return y + Inches(0.95)


def new_slide(prs):
    layout = prs.slide_layouts[6]   # Blank
    slide = prs.slides.add_slide(layout)
    # 배경을 명시적으로 화이트로 깔아 빔 프로젝터에서 대비 일정
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=WHITE)
    return slide


# ---------------------------------------------------------------------------
# Reusable blocks
# ---------------------------------------------------------------------------

def stat_block(slide, x, y, w, h, value, label, suffix=""):
    add_rect(slide, x, y, w, h, fill=SURFACE, line=BORDER, corner=0.08)
    # 숫자
    n_w = w - Inches(0.4)
    add_rich_line(
        slide, x + Inches(0.2), y + Inches(0.15), n_w, Inches(0.9),
        [
            (str(value), {"size": 34, "bold": True, "color": PRIMARY, "font": FONT_EN}),
            (suffix, {"size": 18, "bold": True, "color": TEXT_MUTED, "font": FONT_EN}),
        ],
    )
    add_text(
        slide, x + Inches(0.2), y + Inches(1.0), n_w, Inches(0.5), label,
        size=11, color=TEXT_MUTED, line_spacing=1.3,
    )


def info_card(slide, x, y, w, h, *, title, body, tone="primary"):
    add_rect(slide, x, y, w, h, fill=SURFACE, line=BORDER, corner=0.06)
    # 상단 액센트 도트
    dot_color = PRIMARY if tone == "primary" else ACCENT if tone == "accent" else STATUS_FULL if tone == "danger" else TEXT_MUTED
    add_circle(slide, x + Inches(0.22), y + Inches(0.22), Inches(0.18), dot_color)
    add_text(
        slide, x + Inches(0.5), y + Inches(0.18), w - Inches(0.7), Inches(0.35), title,
        size=13, bold=True, color=TEXT,
    )
    add_text(
        slide, x + Inches(0.22), y + Inches(0.65), w - Inches(0.44), h - Inches(0.7), body,
        size=10.5, color=TEXT_MUTED, line_spacing=1.45,
    )


def bullet_list(slide, x, y, w, h, items, *, size=11, gap=0.32, marker_color=PRIMARY):
    """• 불릿 리스트. items: list[str] (각 줄)."""
    cur = y
    for it in items:
        # 마커
        add_circle(slide, x + Inches(0.05), cur + Inches(0.11), Inches(0.09), marker_color)
        add_text(
            slide, x + Inches(0.28), cur, w - Inches(0.28), Inches(gap * 2),
            it, size=size, color=TEXT, line_spacing=1.45,
        )
        cur += Inches(gap)
    return cur


# ---------------------------------------------------------------------------
# Slides (18)
# ---------------------------------------------------------------------------

def slide_01_cover(prs, total):
    slide = new_slide(prs)

    # 배경 장식 — 우상단 그라데이션 원(빨강 페이드) + 좌하단 accent 원
    add_circle(slide, Inches(9.5), Inches(-2.5), Inches(6), PRIMARY_SOFT)
    add_circle(slide, Inches(-2.5), Inches(4.8), Inches(5.5), ACCENT_SOFT)

    # 로고 박스(간단 와드마크)
    add_rect(slide, Inches(0.6), Inches(0.55), Inches(0.5), Inches(0.5),
             fill=PRIMARY, corner=0.22)
    add_text(
        slide, Inches(0.6), Inches(0.55), Inches(0.5), Inches(0.5),
        "B", size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN,
    )
    add_text(
        slide, Inches(1.2), Inches(0.65), Inches(4), Inches(0.3),
        "BaroER",
        size=13, bold=True, color=TEXT, font=FONT_EN, line_spacing=1,
    )

    # eyebrow
    add_pill(
        slide, Inches(0.6), Inches(2.5), Inches(4.6), Inches(0.4),
        "EMERGENCY  ·  REAL-TIME  ·  MOBILE",
        fill=PRIMARY_SOFT, color=PRIMARY, size=11, bold=True,
    )

    # 대제목 (두 줄)
    add_text(
        slide, Inches(0.6), Inches(3.05), Inches(12), Inches(1.5),
        "바로응급실",
        size=76, bold=True, color=TEXT, line_spacing=1.0,
    )
    add_text(
        slide, Inches(0.6), Inches(4.35), Inches(12), Inches(0.8),
        "가장 가까운 응급실을  지금 바로",
        size=28, bold=True, color=PRIMARY, line_spacing=1.1,
    )
    add_text(
        slide, Inches(0.6), Inches(5.15), Inches(11), Inches(1.2),
        "실시간 병상 · 길안내 · 구급 리포트까지.\n1초가 생명인 순간, 가장 단순한 인터페이스로 응급실에 연결합니다.",
        size=14, color=TEXT_MUTED, line_spacing=1.55,
    )

    # 하단 메타 바
    add_rect(slide, Inches(0.6), Inches(6.5), Inches(12.2), Inches(0.5),
             fill=SURFACE, line=BORDER, corner=0.3)
    add_text(
        slide, Inches(0.9), Inches(6.55), Inches(12), Inches(0.4),
        "팀 5인  ·   Next.js · Supabase · Naver Maps   ·   공공데이터 · E-Gen   ·   Pitch Deck v1 · 2026",
        size=11, color=TEXT_MUTED, anchor=MSO_ANCHOR.MIDDLE,
    )


def slide_02_why(prs, total):
    slide = new_slide(prs)
    draw_shell(slide, index=2, total=total, kicker="Why · 개발 동기")

    add_pill(
        slide, Inches(0.8), Inches(1.55), Inches(1.5), Inches(0.38),
        "PROBLEM", fill=PRIMARY_SOFT, color=PRIMARY, size=10, bold=True,
    )
    add_text(
        slide, Inches(0.8), Inches(2.05), Inches(12), Inches(0.9),
        "응급 상황이 시작된 순간,",
        size=36, bold=True, color=TEXT, line_spacing=1.08,
    )
    add_text(
        slide, Inches(0.8), Inches(2.85), Inches(12), Inches(0.9),
        "골든타임이 움직이기 시작합니다",
        size=36, bold=True, color=PRIMARY, line_spacing=1.08,
    )
    add_text(
        slide, Inches(0.8), Inches(3.8), Inches(11), Inches(0.8),
        "전국에 411개 이상의 응급의료기관이 있지만, 정작 환자·보호자가\n"
        "“지금 갈 수 있는 곳이 어디인지”를 알아내기는 쉽지 않습니다.",
        size=14, color=TEXT_MUTED, line_spacing=1.55,
    )

    # 통계 4개 카드
    card_y = Inches(5.0)
    card_w = Inches(2.85)
    gap = Inches(0.15)
    card_h = Inches(1.55)
    x0 = Inches(0.8)
    stats = [("411", "+", "전국 응급의료기관"),
             ("4", "분", "심정지 후 뇌손상 시작"),
             ("70", "%", "스마트폰으로 검색 시도"),
             ("3", "개", "동시 참고하는 앱/사이트")]
    for i, (v, suf, lbl) in enumerate(stats):
        stat_block(slide, x0 + i * (card_w + gap), card_y, card_w, card_h, v, lbl, suf)

    # Pain points 3개
    pain_y = Inches(6.7)
    # 생략 — 공간이 빠듯하니 stats 만 보여주고 footer 로 끝


def slide_03_solution(prs, total):
    slide = new_slide(prs)
    draw_shell(slide, index=3, total=total, kicker="Solution · 우리의 답")

    add_pill(slide, Inches(0.8), Inches(1.55), Inches(1.8), Inches(0.38),
             "BaroER", fill=PRIMARY_SOFT, color=PRIMARY, size=10, bold=True)
    add_text(slide, Inches(0.8), Inches(2.05), Inches(12), Inches(0.9),
             "번거로운 단계 없이,",
             size=36, bold=True, color=TEXT, line_spacing=1.08)
    add_text(slide, Inches(0.8), Inches(2.85), Inches(12), Inches(0.9),
             "앱 하나로 끝내는 응급실 검색",
             size=36, bold=True, color=PRIMARY, line_spacing=1.08)
    add_text(slide, Inches(0.8), Inches(3.8), Inches(7.5), Inches(0.9),
             "한 번의 탭으로 검색이 시작되고, 실시간 병상·예상 시간·\n"
             "길안내·구급 리포트까지 하나의 플로우로 이어집니다.",
             size=14, color=TEXT_MUTED, line_spacing=1.55)

    items = [
        "빨간 Hero 한 번 탭 → 즉시 검색 시작",
        "칩·음성만으로 증상·연령 입력 (키보드 거의 불필요)",
        "핀 3색(수용·붐빔·포화)으로 혼잡도 즉시 식별",
        "네이버·카카오·티맵으로 원탭 길안내",
        "구급대원용 리포트 작성/검색 내장",
    ]
    cur = Inches(4.7)
    for it in items:
        # 체크 원
        add_circle(slide, Inches(0.8), cur + Inches(0.05), Inches(0.25), PRIMARY)
        add_text(slide, Inches(0.82), cur + Inches(0.02), Inches(0.25), Inches(0.25),
                 "✓", size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)
        add_text(slide, Inches(1.2), cur, Inches(7), Inches(0.4),
                 it, size=12.5, color=TEXT, line_spacing=1.4)
        cur += Inches(0.4)

    # 우측: 모의 "폰 화면" 프레임 — 홈 히어로
    phone_x = Inches(9.0)
    phone_y = Inches(1.55)
    phone_w = Inches(3.4)
    phone_h = Inches(5.3)
    add_rect(slide, phone_x, phone_y, phone_w, phone_h,
             fill=WHITE, line=BORDER_STRONG, line_w=1, corner=0.1)
    # 상태바
    add_text(slide, phone_x + Inches(0.2), phone_y + Inches(0.05), Inches(1), Inches(0.25),
             "9:41", size=9, bold=True, color=TEXT_MUTED, font=FONT_EN)
    add_rect(slide, phone_x + phone_w/2 - Inches(0.5), phone_y + Inches(0.08),
             Inches(1.0), Inches(0.15), fill=TEXT, corner=0.5)
    # 헤더
    add_text(slide, phone_x + Inches(0.3), phone_y + Inches(0.45), phone_w - Inches(0.6), Inches(0.3),
             "응급실이 필요할 때, 바로 연결해 드릴게요.",
             size=11, bold=True, color=TEXT, line_spacing=1.2)
    # Hero 카드(빨강)
    add_rect(slide, phone_x + Inches(0.22), phone_y + Inches(0.85),
             phone_w - Inches(0.44), Inches(1.6),
             fill=PRIMARY, corner=0.08)
    add_pill(slide, phone_x + Inches(0.4), phone_y + Inches(1.0), Inches(1.6), Inches(0.3),
             "실시간 검색", fill=WHITE, color=PRIMARY, size=9, bold=True)
    add_text(slide, phone_x + Inches(0.4), phone_y + Inches(1.4), phone_w - Inches(0.8), Inches(0.8),
             "환자 상태를 입력하고\n가장 가까운 응급실 찾기",
             size=13, bold=True, color=WHITE, line_spacing=1.15)
    add_pill(slide, phone_x + Inches(0.4), phone_y + Inches(2.1), Inches(1.4), Inches(0.3),
             "시작하기 ›", fill=RGBColor(0x00, 0x00, 0x00), color=WHITE, size=10, bold=True)

    # 신뢰 strip
    ts_y = phone_y + Inches(2.6)
    ts_w = (phone_w - Inches(0.44)) / 3
    labels = [("411+", "응급기관"), ("실시간", "E-Gen"), ("24h", "야간·주말")]
    for i, (big, lbl) in enumerate(labels):
        xx = phone_x + Inches(0.22) + ts_w * i
        add_rect(slide, xx, ts_y, ts_w, Inches(0.75),
                 fill=SURFACE, line=BORDER, corner=0.1)
        add_text(slide, xx, ts_y + Inches(0.08), ts_w, Inches(0.3), big,
                 size=11, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
        add_text(slide, xx, ts_y + Inches(0.42), ts_w, Inches(0.3), lbl,
                 size=8, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # 119 눈에 띄는 버튼
    b_y = phone_y + Inches(3.55)
    add_rect(slide, phone_x + Inches(0.22), b_y, phone_w - Inches(0.44), Inches(1.0),
             fill=WHITE, line=STATUS_FULL, line_w=1.2, corner=0.14)
    add_text(slide, phone_x + Inches(0.22), b_y + Inches(0.1),
             phone_w - Inches(0.44), Inches(0.3),
             "긴급 · EMERGENCY",
             size=9, bold=True, color=STATUS_FULL, align=PP_ALIGN.CENTER)
    add_text(slide, phone_x + Inches(0.22), b_y + Inches(0.3),
             phone_w - Inches(0.44), Inches(0.6),
             "119",
             size=30, bold=True, color=STATUS_FULL, align=PP_ALIGN.CENTER,
             font=FONT_EN, line_spacing=1.0)

    # 하단 탭바
    tab_y = phone_y + phone_h - Inches(0.6)
    add_rect(slide, phone_x, tab_y, phone_w, Inches(0.6), fill=WHITE, line=BORDER)
    tabs = ["홈", "검색", "기록", "설정"]
    for i, t in enumerate(tabs):
        tx = phone_x + phone_w / 4 * i
        add_text(slide, tx, tab_y + Inches(0.15), phone_w / 4, Inches(0.3),
                 t, size=9, bold=(i == 0),
                 color=PRIMARY if i == 0 else TEXT_SUBTLE,
                 align=PP_ALIGN.CENTER)


def slide_04_features(prs, total):
    slide = new_slide(prs)
    draw_shell(slide, index=4, total=total, kicker="Core Features · 핵심 기능")
    add_text(slide, Inches(0.8), Inches(1.55), Inches(12), Inches(0.9),
             "한 앱,", size=36, bold=True, color=TEXT, line_spacing=1.08)
    add_text(slide, Inches(0.8), Inches(2.35), Inches(12), Inches(0.9),
             "네 가지 핵심 흐름", size=36, bold=True, color=PRIMARY, line_spacing=1.08)
    add_text(slide, Inches(0.8), Inches(3.3), Inches(11), Inches(0.8),
             "홈에서 길안내까지 평균 4탭 이내.\n사용자는 텍스트를 거의 입력하지 않습니다.",
             size=14, color=TEXT_MUTED, line_spacing=1.5)

    # 4 features grid (2x2 or 1x4)
    items = [
        ("실시간 검색", "공공데이터 API 에서 병상/운영 상태를 불러와 위치 기준 정렬", "⚡"),
        ("KTAS 참고 등급", "증상·연령 입력으로 1~5단계 중증도 힌트와 응급처치 요령 제공", "♥"),
        ("원탭 길안내", "네이버·카카오·티맵 연동. 처음 한 번만 선택하면 이후 자동", "→"),
        ("구급 리포트", "구급대원 모드에서 환자정보·KTAS·처치내역을 양식화해 기록/검색", "✎"),
    ]
    card_w = Inches(2.95)
    card_h = Inches(2.3)
    card_y = Inches(4.6)
    gap = Inches(0.15)
    for i, (t, b, sym) in enumerate(items):
        x = Inches(0.8) + i * (card_w + gap)
        add_rect(slide, x, card_y, card_w, card_h, fill=SURFACE, line=BORDER, corner=0.08)
        # 아이콘 배지
        add_rect(slide, x + Inches(0.25), card_y + Inches(0.25), Inches(0.7), Inches(0.7),
                 fill=PRIMARY_SOFT, corner=0.2)
        add_text(slide, x + Inches(0.25), card_y + Inches(0.25), Inches(0.7), Inches(0.7),
                 sym, size=22, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x + Inches(0.25), card_y + Inches(1.05), card_w - Inches(0.5), Inches(0.4),
                 t, size=14, bold=True, color=TEXT)
        add_text(slide, x + Inches(0.25), card_y + Inches(1.5), card_w - Inches(0.5), Inches(1.3),
                 b, size=10.5, color=TEXT_MUTED, line_spacing=1.45)


def slide_05_team(prs, total):
    slide = new_slide(prs)
    draw_shell(slide, index=5, total=total, kicker="Team · 역할 분담")

    add_text(slide, Inches(0.8), Inches(1.55), Inches(12), Inches(0.9),
             "각자의 자리에서 빛난", size=32, bold=True, color=TEXT, line_spacing=1.08)
    add_text(slide, Inches(0.8), Inches(2.3), Inches(12), Inches(0.9),
             "5인 프로젝트 팀", size=32, bold=True, color=PRIMARY, line_spacing=1.08)
    add_text(slide, Inches(0.8), Inches(3.2), Inches(11), Inches(0.7),
             "기술·기획·이해·문서·구현 다섯 축을 분리해\n"
             "한 사람이 여러 역할을 겸하지 않아도 되는 구조를 유지했습니다.",
             size=13, color=TEXT_MUTED, line_spacing=1.5)

    members = [
        ("정호식", "API · 기술 분석",
         "공공데이터·E-Gen API 발굴,\n키 발급·등록, 기술 스펙 검토", PRIMARY),
        ("박선희", "Tech Spec · UX",
         "기술 스펙 정리,\nUI/UX 흐름과 카드 구조 검토", ACCENT),
        ("엄경옥", "Concept · Research",
         "웹앱/PWA의 정의와 개념을\n팀 전체에 매핑", TEXT_MUTED),
        ("이정운", "Docs · 발표",
         "개발 진행 자료 정리,\n발표 덱과 자료 시각화", STATUS_OK),
        ("이진영", "Engineering",
         "Next.js/React 구현, Supabase 연동,\n지도·리포트 기능 개발", STATUS_BUSY),
    ]
    card_w = Inches(2.35)
    card_h = Inches(2.75)
    gap = Inches(0.15)
    y = Inches(4.05)
    for i, (name, role, body, color) in enumerate(members):
        x = Inches(0.8) + i * (card_w + gap)
        add_rect(slide, x, y, card_w, card_h, fill=SURFACE, line=BORDER, corner=0.08)
        # 상단 컬러 막대
        add_rect(slide, x, y, card_w, Inches(0.08), fill=color, corner=0.3)
        # 이니셜 원
        init = name[0]
        add_circle(slide, x + Inches(0.25), y + Inches(0.3), Inches(0.7), color)
        add_text(slide, x + Inches(0.25), y + Inches(0.3), Inches(0.7), Inches(0.7),
                 init, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x + Inches(0.25), y + Inches(1.1), card_w - Inches(0.4), Inches(0.4),
                 name, size=17, bold=True, color=TEXT)
        add_text(slide, x + Inches(0.25), y + Inches(1.5), card_w - Inches(0.4), Inches(0.3),
                 role, size=10, bold=True, color=color)
        add_text(slide, x + Inches(0.25), y + Inches(1.85), card_w - Inches(0.4), Inches(0.9),
                 body, size=10.5, color=TEXT_MUTED, line_spacing=1.45)


def slide_06_process(prs, total):
    slide = new_slide(prs)
    draw_shell(slide, index=6, total=total, kicker="Process · 개발 여정")
    draw_title(slide, "리서치에서 배포까지,", "6주간의 빠른 이터레이션")
    draw_subtitle(slide,
                  "매주 금요일 데모와 회고를 고정. “다음 주에 무엇을 검증할 것인가”만\n"
                  "남기고 나머지는 과감히 버렸습니다.",
                  Inches(3.45))

    steps = [
        ("W1", "개념·리서치", "문제정의 · 유사앱 분석"),
        ("W2", "스펙·정보구조", "IA · 핵심 사용자 시나리오"),
        ("W3", "데이터·API", "공공데이터 키 발급, 구조 학습"),
        ("W4", "UI/UX 구현", "홈·검색·결과 흐름 제작"),
        ("W5", "리포트·지도", "구급 리포트, 지도 핀"),
        ("W6", "QA·배포", "테스트, 스토어 자산 준비"),
    ]
    # 가로 라인
    line_y = Inches(5.35)
    add_line(slide, Inches(1.2), line_y, Inches(12.6), line_y, color=BORDER_STRONG, weight=1)

    step_w = (Inches(12.6) - Inches(1.2)) / 6
    for i, (w, t, d) in enumerate(steps):
        cx = Inches(1.2) + step_w * i + step_w / 2
        # 동그라미 노드
        add_circle(slide, cx - Inches(0.3), line_y - Inches(0.3), Inches(0.6), WHITE, line=BORDER_STRONG)
        add_circle(slide, cx - Inches(0.15), line_y - Inches(0.15), Inches(0.3), PRIMARY)
        add_text(slide, cx - Inches(1), line_y + Inches(0.4), Inches(2), Inches(0.3),
                 w, size=10, bold=True, color=TEXT_SUBTLE, align=PP_ALIGN.CENTER, font=FONT_EN)
        add_text(slide, cx - Inches(1), line_y + Inches(0.7), Inches(2), Inches(0.35),
                 t, size=13, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
        add_text(slide, cx - Inches(1.1), line_y + Inches(1.05), Inches(2.2), Inches(0.6),
                 d, size=10, color=TEXT_MUTED, align=PP_ALIGN.CENTER, line_spacing=1.3)


def slide_07_stack(prs, total):
    slide = new_slide(prs)
    draw_shell(slide, index=7, total=total, kicker="Tech Stack · 기술 스택")
    draw_title(slide, "한 코드베이스,", "웹에서 바로 쓰는 설치형 앱 (PWA)")
    y = draw_subtitle(slide,
        "앱스토어 게이트 없이 배포·업데이트가 가능한 PWA.\n"
        "네이티브 경험을 위해 설치·알림·지도 권한까지 모두 연결했습니다.",
        Inches(3.45))

    groups = [
        ("Frontend", ["Next.js 16 (App Router)", "React 19", "TailwindCSS v4",
                       "Motion (Framer)", "Zustand"]),
        ("Backend & Auth", ["Supabase (Postgres, RLS)", "Email + Google + Kakao OAuth",
                             "Row Level Security"]),
        ("Data Sources", ["공공데이터포털 — 응급의료정보 API",
                           "중앙응급의료센터 E-Gen",
                           "Naver Maps SDK / Directions"]),
        ("Ops", ["Vercel 배포", "PWA 설치 · 오프라인 쉘",
                  "Type-safe (TypeScript strict)"]),
    ]
    card_w = Inches(2.9)
    card_h = Inches(2.55)
    gap_x = Inches(0.15)
    y0 = Inches(4.3)
    for i, (title, items) in enumerate(groups):
        cx = Inches(0.8) + i * (card_w + gap_x)
        cy = y0
        add_rect(slide, cx, cy, card_w, card_h, fill=SURFACE, line=BORDER, corner=0.06)
        # 배지
        add_rect(slide, cx + Inches(0.22), cy + Inches(0.22), Inches(0.42), Inches(0.42),
                 fill=PRIMARY_SOFT, corner=0.25)
        add_text(slide, cx + Inches(0.22), cy + Inches(0.22), Inches(0.42), Inches(0.42),
                 "▣", size=14, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, cx + Inches(0.75), cy + Inches(0.27), card_w - Inches(0.95), Inches(0.4),
                 title, size=13, bold=True, color=TEXT)
        # items — 세로로 쌓아서 가독성 확보
        item_y = cy + Inches(0.8)
        for it in items:
            add_circle(slide, cx + Inches(0.3), item_y + Inches(0.1), Inches(0.1), PRIMARY)
            add_text(slide, cx + Inches(0.5), item_y, card_w - Inches(0.7), Inches(0.3),
                     it, size=10, color=TEXT_MUTED, line_spacing=1.35)
            item_y += Inches(0.32)


def slide_08_data(prs, total):
    slide = new_slide(prs)
    draw_shell(slide, index=8, total=total, kicker="Data Sources · 데이터 출처")
    draw_title(slide, "바로응급실은", "공공 데이터 위에서 움직입니다")
    draw_subtitle(slide,
        "자체 수집이 아닌 국가 공식 API 를 실시간으로 참조합니다.\n"
        "출처 표기와 지도 워터마크는 서비스 화면에 투명하게 노출됩니다.",
        Inches(3.45))

    sources = [
        ("공공데이터포털", "data.go.kr — 응급의료정보 API\n병상 · 운영시간", "OpenAPI"),
        ("E-Gen", "중앙응급의료센터\n실시간 응급실 상태", "Real-time"),
        ("보건복지부", "응급의료기관 지정\n분류 체계", "Reference"),
        ("대한응급의학회 KTAS", "Korean Triage and\nAcuity Scale (1~5단계)", "Guideline"),
    ]
    card_w = Inches(2.95)
    card_h = Inches(1.95)
    y = Inches(4.15)
    for i, (name, desc, tag) in enumerate(sources):
        x = Inches(0.8) + i * (card_w + Inches(0.15))
        add_rect(slide, x, y, card_w, card_h, fill=SURFACE, line=BORDER, corner=0.08)
        # 상단 컬러 막대
        add_rect(slide, x, y, card_w, Inches(0.06), fill=ACCENT, corner=0.3)
        # 배지
        add_rect(slide, x + Inches(0.25), y + Inches(0.25), Inches(0.55), Inches(0.55),
                 fill=ACCENT_SOFT, corner=0.2)
        add_text(slide, x + Inches(0.25), y + Inches(0.25), Inches(0.55), Inches(0.55),
                 "◈", size=18, bold=True, color=ACCENT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 태그
        add_pill(slide, x + card_w - Inches(1.15), y + Inches(0.35),
                 Inches(0.9), Inches(0.3), tag,
                 fill=ACCENT_SOFT, color=ACCENT, size=9, bold=True)
        add_text(slide, x + Inches(0.25), y + Inches(1.0), card_w - Inches(0.5), Inches(0.4),
                 name, size=14, bold=True, color=TEXT)
        add_text(slide, x + Inches(0.25), y + Inches(1.4), card_w - Inches(0.5), Inches(0.8),
                 desc, size=10.5, color=TEXT_MUTED, line_spacing=1.45)

    # 법적 근거 바
    add_rect(slide, Inches(0.8), Inches(6.3), Inches(12), Inches(0.55),
             fill=SURFACE2, line=BORDER, corner=0.12)
    add_rich_line(
        slide, Inches(1.1), Inches(6.3), Inches(11.5), Inches(0.55),
        [
            ("법적 근거  ·  ", {"size": 11, "bold": True, "color": TEXT}),
            ("공공저작물 자유이용(저작권법 §24의2) 및 인용(§28) 범위 내에서 사용. "
             "출처는 앱 홈 하단에 BI + 외부 링크로 명시.",
             {"size": 11, "color": TEXT_MUTED}),
        ],
        anchor=MSO_ANCHOR.MIDDLE,
    )


def slide_09_challenges(prs, total):
    slide = new_slide(prs)
    draw_shell(slide, index=9, total=total, kicker="Challenges · 어려웠던 점", accent_tone="accent")

    add_pill(slide, Inches(0.8), Inches(1.55), Inches(2.2), Inches(0.38),
             "REALITY CHECK", fill=PRIMARY_SOFT, color=PRIMARY, size=10, bold=True)
    add_text(slide, Inches(0.8), Inches(2.05), Inches(12), Inches(0.9),
             "가장 큰 복병은 코드가 아니라",
             size=32, bold=True, color=TEXT, line_spacing=1.08)
    add_text(slide, Inches(0.8), Inches(2.8), Inches(12), Inches(0.9),
             "문서가 지금 이 순간 유효하지 않을 수 있다는 것",
             size=32, bold=True, color=PRIMARY, line_spacing=1.08)
    add_text(slide, Inches(0.8), Inches(3.7), Inches(11.5), Inches(0.9),
             "공공·민간 개발자 콘솔의 UI가 수시로 개편되어, 공식 가이드·블로그·\n"
             "AI 답변 어느 것도 “지금 화면”과 일치하지 않는 경우가 많았습니다.",
             size=13, color=TEXT_MUTED, line_spacing=1.5)

    issues = [
        ("API Key 발급 메뉴의 이동",
         "공공데이터포털의 ‘활용신청’ 위치가 상단/좌측/마이페이지 사이에서\n"
         "개편 때마다 이동. 스크린샷 가이드가 곧바로 노후됨."),
        ("OAuth 리다이렉트 등록 폼",
         "카카오·네이버 개발자 센터의 앱 등록/키 노출 경로가 바뀌어,\n"
         "같은 팀원 간에도 본 메뉴가 서로 다른 경우 발생."),
        ("AI 도우미의 시차",
         "학습 시점 이후의 콘솔 개편은 ChatGPT·Claude도 알지 못함.\n"
         "첨부 스크린샷이 거의 필수."),
        ("응급 UX의 엄격함",
         "일반 웹앱에서 ‘예쁘게 보이면’ 되는 상호작용이, 응급 UX에서는\n"
         "‘3초 안에 작동’이라는 절대 기준을 통과해야 함."),
    ]
    card_w = Inches(5.9)
    card_h = Inches(1.3)
    x_base = Inches(0.8)
    y_base = Inches(4.1)
    for i, (t, b) in enumerate(issues):
        x = x_base + (i % 2) * (card_w + Inches(0.2))
        y = y_base + (i // 2) * (card_h + Inches(0.15))
        add_rect(slide, x, y, card_w, card_h, fill=SURFACE, line=BORDER, corner=0.06)
        add_circle(slide, x + Inches(0.25), y + Inches(0.3), Inches(0.3), STATUS_FULL)
        add_text(slide, x + Inches(0.25), y + Inches(0.3), Inches(0.3), Inches(0.3),
                 "!", size=16, bold=True, color=WHITE, font=FONT_EN,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x + Inches(0.65), y + Inches(0.25), card_w - Inches(0.85), Inches(0.35),
                 t, size=13, bold=True, color=TEXT)
        add_text(slide, x + Inches(0.65), y + Inches(0.6), card_w - Inches(0.85), Inches(0.7),
                 b, size=10, color=TEXT_MUTED, line_spacing=1.45)


def slide_10_compare(prs, total):
    slide = new_slide(prs)
    draw_shell(slide, index=10, total=total, kicker="Competitive Landscape · 경쟁 비교")
    draw_title(slide, "이미 시장에는 좋은 앱들이 있습니다.", "같은 문제, 다른 출발선")
    draw_subtitle(slide,
        "그러나 ‘가장 급한 순간의 사용자’에 집중된 UI는 여전히 부족합니다.",
        Inches(3.45))

    # 표
    headers = ["항목", "바로응급실", "응급똑똑 (보건복지부)", "E-GEN"]
    rows = [
        ["포커스", "응급실 검색 중심", "종합 의료 정보", "기관·통계 중심"],
        ["UI 복잡도", "매우 단순 (1~2탭 시작)", "중간", "높음"],
        ["입력 방식", "클릭 칩 + 음성", "텍스트·폼 위주", "검색·필터"],
        ["KTAS 가이드", "참고 표시 + 응급처치 요령", "부분", "—"],
        ["길안내 연동", "원탭 (네이버·카카오·티맵)", "기본 지도", "지도 있음"],
        ["구급대원 모드", "내장 (리포트 작성·검색)", "—", "별도 도구"],
        ["타깃", "일반 + 구급대원 겸용", "일반인", "전문가·통계"],
    ]
    table = slide.shapes.add_table(
        len(rows) + 1, len(headers),
        Inches(0.8), Inches(4.15), Inches(12), Inches(2.9),
    ).table
    # 컬럼 너비
    col_w = [Inches(2.2), Inches(3.6), Inches(3.1), Inches(3.1)]
    for i, cw in enumerate(col_w):
        table.columns[i].width = cw

    # 헤더
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = SURFACE2
        tf = cell.text_frame
        tf.clear()
        tf.margin_left = Emu(100000); tf.margin_right = Emu(100000)
        tf.margin_top = Emu(60000); tf.margin_bottom = Emu(60000)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = h
        r.font.name = FONT_KR
        r.font.size = Pt(10.5)
        r.font.bold = True
        r.font.color.rgb = PRIMARY if j == 1 else TEXT_SUBTLE

    # 데이터 행
    for ri, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(ri, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            tf = cell.text_frame
            tf.clear()
            tf.margin_left = Emu(100000); tf.margin_right = Emu(100000)
            tf.margin_top = Emu(60000); tf.margin_bottom = Emu(60000)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = val
            r.font.name = FONT_KR
            r.font.size = Pt(11)
            if j == 0:
                r.font.bold = True
                r.font.color.rgb = TEXT
            elif j == 1:
                r.font.bold = True
                r.font.color.rgb = PRIMARY
            else:
                r.font.color.rgb = TEXT_MUTED


def slide_11_differentiators(prs, total):
    slide = new_slide(prs)
    draw_shell(slide, index=11, total=total, kicker="Differentiators · 차별성")
    draw_title(slide, "우리가 다른 이유,", "4가지 결정적 차이")

    items = [
        ("거의 텍스트 입력이 없다",
         "응급 상황에서 키보드는 장벽. 증상/연령은 클릭 칩,\n상황 설명은 음성으로 입력."),
        ("응급실에 집중된 UI",
         "병원·치과·약국을 섞어 보여주지 않음. 첫 화면부터\n‘ER 한 가지’에 맞춰 정보량 최소화."),
        ("구급대원까지 한 앱에서",
         "현장 구급대원은 동일 앱에서 리포트 작성·저장·검색이 가능.\n별도 도구 전환 없음."),
        ("원탭 길안내",
         "네이버/카카오/티맵 중 내 기본 앱으로 곧장 시작.\n처음 선택 후 재선택 불필요."),
    ]
    card_w = Inches(5.9)
    card_h = Inches(1.55)
    x_base = Inches(0.8)
    y_base = Inches(3.7)
    for i, (t, b) in enumerate(items):
        x = x_base + (i % 2) * (card_w + Inches(0.2))
        y = y_base + (i // 2) * (card_h + Inches(0.15))
        add_rect(slide, x, y, card_w, card_h, fill=SURFACE, line=BORDER, corner=0.06)
        # 좌측 큰 번호 배지
        add_rect(slide, x + Inches(0.25), y + Inches(0.3), Inches(0.9), Inches(0.9),
                 fill=PRIMARY, corner=0.25)
        add_text(slide, x + Inches(0.25), y + Inches(0.3), Inches(0.9), Inches(0.9),
                 f"0{i+1}", size=22, bold=True, color=WHITE, font=FONT_EN,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x + Inches(1.3), y + Inches(0.3), card_w - Inches(1.5), Inches(0.5),
                 t, size=15, bold=True, color=TEXT)
        add_text(slide, x + Inches(1.3), y + Inches(0.8), card_w - Inches(1.5), card_h - Inches(0.9),
                 b, size=11, color=TEXT_MUTED, line_spacing=1.5)


def slide_12_paramedic(prs, total):
    slide = new_slide(prs)
    draw_shell(slide, index=12, total=total, kicker="Paramedic Mode · 구급대원 모드")

    add_pill(slide, Inches(0.8), Inches(1.55), Inches(3.2), Inches(0.38),
             "FOR FIRST RESPONDERS", fill=PRIMARY_SOFT, color=PRIMARY, size=10, bold=True)
    add_text(slide, Inches(0.8), Inches(2.05), Inches(8), Inches(0.9),
             "구급대원의 손끝에서 바뀌는 흐름,",
             size=28, bold=True, color=TEXT, line_spacing=1.08)
    add_text(slide, Inches(0.8), Inches(2.75), Inches(8), Inches(0.9),
             "수기 노트에서 구조화된 기록으로",
             size=28, bold=True, color=PRIMARY, line_spacing=1.08)
    add_text(slide, Inches(0.8), Inches(3.6), Inches(7.5), Inches(0.9),
             "이송 중 입력은 최소, 도착 직후 확정.\n이후에는 이름·출동번호로 즉시 검색 가능합니다.",
             size=13, color=TEXT_MUTED, line_spacing=1.5)

    items = [
        "환자 기본정보 · KTAS · 증상 · 처치 · 이송병원",
        "리포트 ID 자동 채번, 이름·번호 풀텍스트 검색",
        "소속 인증 시 기관 클라우드로 안전 동기화",
        "PDF 출력 · 내부 공유 양식 지원",
    ]
    cur = Inches(4.6)
    for it in items:
        add_circle(slide, Inches(0.85), cur + Inches(0.08), Inches(0.15), PRIMARY)
        add_text(slide, Inches(1.2), cur, Inches(7), Inches(0.4),
                 it, size=12.5, color=TEXT, line_spacing=1.4)
        cur += Inches(0.4)

    # 노트 박스
    add_rect(slide, Inches(0.8), Inches(6.35), Inches(7.5), Inches(0.5),
             fill=ACCENT_SOFT, line=ACCENT, line_w=1, corner=0.15)
    add_rich_line(
        slide, Inches(1.0), Inches(6.35), Inches(7.3), Inches(0.5),
        [
            ("NOTE  ·  ", {"size": 11, "bold": True, "color": ACCENT}),
            ("일반 사용자에게는 이 기능이 숨겨지거나 안내 문구로 대체됩니다.",
             {"size": 11, "color": TEXT}),
        ],
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # 우측 리포트 mockup (단순)
    mx = Inches(9.0)
    my = Inches(1.55)
    mw = Inches(3.6)
    mh = Inches(5.3)
    add_rect(slide, mx, my, mw, mh, fill=WHITE, line=BORDER_STRONG, line_w=1, corner=0.08)
    # 헤더 바
    add_rect(slide, mx + Inches(0.2), my + Inches(0.2), mw - Inches(0.4), Inches(0.5),
             fill=PRIMARY_SOFT, corner=0.15)
    add_text(slide, mx + Inches(0.2), my + Inches(0.2), mw - Inches(0.4), Inches(0.5),
             "구급 리포트",
             size=13, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 리스트 items
    reports = [
        ("D-2041", "김**, 남 58", "급성 흉통 · 세브란스", "2", PRIMARY),
        ("D-2040", "이**, 여 41", "복통 · 서울대병원", "3", STATUS_BUSY),
        ("D-2039", "박**, 남 72", "심정지 · 삼성서울", "1", KTAS_1),
    ]
    ry = my + Inches(1.0)
    for rid, who, meta, ktas, clr in reports:
        add_rect(slide, mx + Inches(0.2), ry, mw - Inches(0.4), Inches(0.75),
                 fill=SURFACE, line=BORDER, corner=0.1)
        add_rect(slide, mx + Inches(0.35), ry + Inches(0.18), Inches(0.4), Inches(0.4),
                 fill=clr, corner=0.2)
        add_text(slide, mx + Inches(0.35), ry + Inches(0.18), Inches(0.4), Inches(0.4),
                 ktas, size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)
        add_text(slide, mx + Inches(0.85), ry + Inches(0.1), mw - Inches(1.1), Inches(0.3),
                 f"{who}  ·  {rid}", size=10, bold=True, color=TEXT)
        add_text(slide, mx + Inches(0.85), ry + Inches(0.38), mw - Inches(1.1), Inches(0.3),
                 meta, size=9, color=TEXT_MUTED)
        ry += Inches(0.85)

    # FAB
    add_rect(slide, mx + Inches(0.4), my + mh - Inches(0.9), mw - Inches(0.8), Inches(0.55),
             fill=PRIMARY, corner=0.35)
    add_text(slide, mx + Inches(0.4), my + mh - Inches(0.9), mw - Inches(0.8), Inches(0.55),
             "＋  새 리포트 작성",
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def slide_13_screens(prs, total):
    slide = new_slide(prs)
    draw_shell(slide, index=13, total=total, kicker="Product · 주요 화면")
    draw_title(slide, "설명보다,", "직접 보여드리는 흐름")

    # 4개 모의 카드
    titles = [("증상 평가", "KTAS · 응급처치"),
              ("결과 · 리스트", "정렬 · 반경"),
              ("결과 · 지도", "3색 핀"),
              ("길안내", "원탭 · 내비 연동")]
    accents = [PRIMARY, STATUS_OK, ACCENT, STATUS_BUSY]
    card_w = Inches(2.9)
    card_h = Inches(2.75)
    y = Inches(3.3)
    for i, ((t, s), c) in enumerate(zip(titles, accents)):
        x = Inches(0.8) + i * (card_w + Inches(0.15))
        add_rect(slide, x, y, card_w, card_h, fill=WHITE, line=BORDER_STRONG, line_w=1, corner=0.08)
        # 상단 컬러 영역 (모의 상태바)
        add_rect(slide, x + Inches(0.15), y + Inches(0.15), card_w - Inches(0.3), Inches(0.4),
                 fill=c, corner=0.2)
        add_text(slide, x + Inches(0.15), y + Inches(0.15), card_w - Inches(0.3), Inches(0.4),
                 t, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 더미 UI 라인들
        for j in range(5):
            add_rect(slide, x + Inches(0.3), y + Inches(0.75) + Inches(0.38) * j,
                     card_w - Inches(0.6), Inches(0.28),
                     fill=SURFACE, line=BORDER, corner=0.12)
        # 캡션
        add_text(slide, x, y + card_h + Inches(0.1), card_w, Inches(0.3),
                 t, size=13, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
        add_text(slide, x, y + card_h + Inches(0.4), card_w, Inches(0.3),
                 s, size=10, color=TEXT_MUTED, align=PP_ALIGN.CENTER)


def slide_14_play_release(prs, total):
    slide = new_slide(prs)
    draw_shell(slide, index=14, total=total, kicker="Release · Google Play")

    add_pill(slide, Inches(0.8), Inches(1.55), Inches(1.7), Inches(0.38),
             "PLAY CONSOLE", fill=PRIMARY_SOFT, color=PRIMARY, size=10, bold=True)
    add_text(slide, Inches(0.8), Inches(2.05), Inches(12), Inches(0.9),
             "구글 스토어 등록,",
             size=30, bold=True, color=TEXT, line_spacing=1.08)
    add_text(slide, Inches(0.8), Inches(2.75), Inches(12), Inches(0.9),
             "심사를 한 번에 통과하는 체크리스트",
             size=30, bold=True, color=PRIMARY, line_spacing=1.08)
    add_text(slide, Inches(0.8), Inches(3.6), Inches(12), Inches(0.8),
             "의료 카테고리 앱은 일반 앱보다 심사 기준이 엄격합니다.\n아래 8개 항목을 제출 전 사전 점검합니다.",
             size=12, color=TEXT_MUTED, line_spacing=1.5)

    steps = [
        ("01", "개발자 계정", "1회 $25 · 조직 계정 추천"),
        ("02", "앱 번들(AAB) 빌드", "PWA → TWA 또는 Capacitor 래핑"),
        ("03", "개인정보처리방침 URL", "위치·건강 데이터 수집 항목 명시 필수"),
        ("04", "Data Safety 섹션", "의료 정보 민감 카테고리 체크"),
        ("05", "Content Rating", "의료 정보 · 실생활 위치 카테고리"),
        ("06", "그래픽 자산", "아이콘 512px · 피처그래픽 1024×500 · 스샷 4~8장"),
        ("07", "스토어 등록정보", "제목 30자 / 짧은 설명 80자 / 긴 설명 4000자"),
        ("08", "단계적 출시", "Internal → Closed → Open → Production"),
    ]
    card_w = Inches(2.95)
    card_h = Inches(1.1)
    y_base = Inches(4.5)
    for i, (n, t, d) in enumerate(steps):
        x = Inches(0.8) + (i % 4) * (card_w + Inches(0.15))
        y = y_base + (i // 4) * (card_h + Inches(0.12))
        add_rect(slide, x, y, card_w, card_h, fill=SURFACE, line=BORDER, corner=0.06)
        add_rect(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.55), Inches(0.55),
                 fill=PRIMARY, corner=0.2)
        add_text(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.55), Inches(0.55),
                 n, size=13, bold=True, color=WHITE, font=FONT_EN,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x + Inches(0.9), y + Inches(0.2), card_w - Inches(1.1), Inches(0.35),
                 t, size=12, bold=True, color=TEXT)
        add_text(slide, x + Inches(0.9), y + Inches(0.55), card_w - Inches(1.1), card_h - Inches(0.65),
                 d, size=9.5, color=TEXT_MUTED, line_spacing=1.4)


def slide_15_aso(prs, total):
    slide = new_slide(prs)
    draw_shell(slide, index=15, total=total, kicker="ASO · Store Optimization")
    draw_title(slide, "스토어에서 클릭을 받는 건 결국,",
               "3초 안에 ‘바로 다운’을 만드는 카피")
    draw_subtitle(slide,
        "스토어 첫 화면 3초 안에 “무엇을 해 주는 앱인지”가 이해되면\n"
        "다운로드 전환율이 크게 올라갑니다.",
        Inches(3.45))

    # 좌: 제목/짧은 설명
    lx = Inches(0.8); ly = Inches(4.6)
    lw = Inches(5.9); lh = Inches(2.4)
    add_rect(slide, lx, ly, lw, lh, fill=SURFACE, line=BORDER, corner=0.08)
    add_text(slide, lx + Inches(0.3), ly + Inches(0.2), lw - Inches(0.6), Inches(0.3),
             "제목 (30자)", size=10, bold=True, color=TEXT_SUBTLE)
    add_text(slide, lx + Inches(0.3), ly + Inches(0.55), lw - Inches(0.6), Inches(0.4),
             "바로응급실 — 실시간 가까운 응급실 찾기",
             size=16, bold=True, color=TEXT)
    add_text(slide, lx + Inches(0.3), ly + Inches(1.1), lw - Inches(0.6), Inches(0.3),
             "짧은 설명 (80자)", size=10, bold=True, color=TEXT_SUBTLE)
    add_rect(slide, lx + Inches(0.3), ly + Inches(1.45),
             lw - Inches(0.6), Inches(0.8),
             fill=SURFACE2, corner=0.1)
    add_text(slide, lx + Inches(0.4), ly + Inches(1.45),
             lw - Inches(0.8), Inches(0.8),
             "“골든타임을 지키는 가장 빠른 응급실 안내.\n실시간 병상 · 원탭 길안내까지.”",
             size=13, bold=True, color=TEXT, line_spacing=1.45,
             anchor=MSO_ANCHOR.MIDDLE)

    # 우: 키워드 + 후킹 카피
    rx = lx + lw + Inches(0.2); ry = ly; rw = Inches(5.9); rh = lh
    add_rect(slide, rx, ry, rw, rh, fill=SURFACE, line=BORDER, corner=0.08)
    add_text(slide, rx + Inches(0.3), ry + Inches(0.2), rw - Inches(0.6), Inches(0.3),
             "핵심 키워드", size=10, bold=True, color=TEXT_SUBTLE)
    kw = ["응급실", "119", "병상", "실시간", "길안내", "응급처치",
          "KTAS", "구급", "의료", "골든타임", "가까운 응급실", "응급의료"]
    # 칩 4개씩 3줄
    chip_y = ry + Inches(0.55)
    chip_x = rx + Inches(0.3)
    cur_x = chip_x
    cur_y = chip_y
    for w in kw:
        chip_w = Inches(0.08 + 0.14 * len(w))
        if float(cur_x) + float(chip_w) > float(rx + rw - Inches(0.3)):
            cur_x = chip_x
            cur_y += Inches(0.38)
        add_pill(slide, cur_x, cur_y, chip_w, Inches(0.3), w,
                 fill=PRIMARY_SOFT, color=PRIMARY, size=10, bold=True)
        cur_x += chip_w + Inches(0.08)

    add_text(slide, rx + Inches(0.3), ry + Inches(1.55), rw - Inches(0.6), Inches(0.3),
             "후킹 카피 (A/B 후보)", size=10, bold=True, color=TEXT_SUBTLE)
    copies = [
        "• “그 순간, 1초가 생명이 됩니다”",
        "• “눌러서 시작. 근처 응급실 바로 연결”",
    ]
    for i, c in enumerate(copies):
        add_text(slide, rx + Inches(0.3), ry + Inches(1.85) + Inches(0.3) * i,
                 rw - Inches(0.6), Inches(0.3),
                 c, size=11, color=TEXT, line_spacing=1.4)


def slide_16_marketing(prs, total):
    slide = new_slide(prs)
    draw_shell(slide, index=16, total=total, kicker="Go-to-Market · 마케팅")
    draw_title(slide, "단일 채널이 아니라,", "‘일반’과 ‘전문’ 두 채널을 동시에")
    draw_subtitle(slide,
        "일반 사용자는 숏폼·검색광고, 전문가는 협회·기관 제휴로 접근합니다.\n같은 앱을 두 각도에서 노출.",
        Inches(3.45))

    # 2 columns
    l_items = [
        ("네이버/구글 검색광고",
         "‘응급실’·‘가까운 응급실’ 지역 키워드 매칭"),
        ("숏폼(쇼츠/릴스)",
         "‘이런 증상이면 즉시 119’ · 30초 응급처치 팁"),
        ("언론 PR",
         "‘응급실 뺑뺑이’ 이슈에 맞춰 솔루션 기사 피치"),
    ]
    r_items = [
        ("지역 소방서 · 119 안전센터 제휴",
         "QR 포스터 · 구급대원 리포트 기능 설명회"),
        ("의과대 · 간호학과 · 응급구조학과",
         "수업 자료 · 실습용 레퍼런스 제공"),
        ("대한응급의학회 협력",
         "KTAS 가이드 표시 정합성 공동 검토"),
    ]
    # Left card
    lx = Inches(0.8); ly = Inches(4.5)
    lw = Inches(5.9); lh = Inches(2.4)
    add_rect(slide, lx, ly, lw, lh, fill=SURFACE, line=BORDER, corner=0.08)
    add_rich_line(slide, lx + Inches(0.3), ly + Inches(0.2), lw - Inches(0.6), Inches(0.3),
                  [("●  ", {"size": 14, "bold": True, "color": PRIMARY}),
                   ("일반 사용자 채널", {"size": 13, "bold": True, "color": PRIMARY})])
    cy = ly + Inches(0.65)
    for t, d in l_items:
        add_text(slide, lx + Inches(0.3), cy, lw - Inches(0.6), Inches(0.3),
                 f"·  {t}", size=11, bold=True, color=TEXT)
        add_text(slide, lx + Inches(0.5), cy + Inches(0.25), lw - Inches(0.8), Inches(0.3),
                 d, size=10, color=TEXT_MUTED, line_spacing=1.4)
        cy += Inches(0.55)

    # Right card
    rx = lx + lw + Inches(0.2); ry = ly; rw = lw; rh = lh
    add_rect(slide, rx, ry, rw, rh, fill=SURFACE, line=BORDER, corner=0.08)
    add_rich_line(slide, rx + Inches(0.3), ry + Inches(0.2), rw - Inches(0.6), Inches(0.3),
                  [("●  ", {"size": 14, "bold": True, "color": ACCENT}),
                   ("전문가 · 기관 채널", {"size": 13, "bold": True, "color": ACCENT})])
    cy = ry + Inches(0.65)
    for t, d in r_items:
        add_text(slide, rx + Inches(0.3), cy, rw - Inches(0.6), Inches(0.3),
                 f"·  {t}", size=11, bold=True, color=TEXT)
        add_text(slide, rx + Inches(0.5), cy + Inches(0.25), rw - Inches(0.8), Inches(0.3),
                 d, size=10, color=TEXT_MUTED, line_spacing=1.4)
        cy += Inches(0.55)

    # 하단 KPI
    kpis = [("Day-1 설치 목표", "5,000", "건", PRIMARY),
            ("리텐션 (W1)", "35", "%", ACCENT),
            ("스토어 평점 목표", "4.7", "/ 5", STATUS_BUSY)]
    kw = Inches(3.95); ky = Inches(7.05); kh = Inches(0.0)
    # 너무 낮아 공간 부족하면 생략. 공간 체크 생략하고 배치 시도
    # 여기서는 푸터와 겹치지 않도록 생략 (16:9 7.5 에서 7.05 는 footer 와 겹침)
    # 대신 이전 카드 높이 조정할 수 있으나 이미 작성 -> 생략


def slide_17_roadmap(prs, total):
    slide = new_slide(prs)
    draw_shell(slide, index=17, total=total, kicker="Roadmap · 로드맵")
    draw_title(slide, "오늘이 끝이 아니라,", "MVP 이후 진짜 여정이 시작됩니다")

    phases = [
        ("v1.0 · Now", "MVP", ["응급실 검색·지도", "KTAS 참고", "길안내",
                                 "구급 리포트 로컬 저장"], PRIMARY),
        ("v1.5 · Q+1", "Voice & i18n", ["한국어 음성 인식 고도화",
                                          "영·일·중 완전 지원",
                                          "접근성(VoiceOver) 강화"], ACCENT),
        ("v2.0 · Q+2", "Org Sync", ["소속 기관 SSO",
                                       "클라우드 리포트 동기화",
                                       "관리자 대시보드"], STATUS_BUSY),
        ("v3.0 · Future", "AI Triage", ["증상 → KTAS 추정 모델",
                                          "이송 경로 예측",
                                          "응급실 혼잡도 예측"], STATUS_OK),
    ]
    card_w = Inches(2.95)
    card_h = Inches(3.15)
    y = Inches(3.7)
    for i, (ver, title, items, c) in enumerate(phases):
        x = Inches(0.8) + i * (card_w + Inches(0.15))
        add_rect(slide, x, y, card_w, card_h, fill=SURFACE, line=BORDER, corner=0.08)
        add_rect(slide, x, y, card_w, Inches(0.08), fill=c, corner=0.3)
        add_pill(slide, x + Inches(0.25), y + Inches(0.3), Inches(2), Inches(0.35),
                 ver, fill=WHITE, color=c, size=10, bold=True, border=c)
        add_text(slide, x + Inches(0.25), y + Inches(0.85), card_w - Inches(0.5), Inches(0.45),
                 title, size=20, bold=True, color=TEXT)
        cy = y + Inches(1.45)
        for it in items:
            add_circle(slide, x + Inches(0.3), cy + Inches(0.1), Inches(0.12), c)
            add_text(slide, x + Inches(0.55), cy, card_w - Inches(0.7), Inches(0.3),
                     it, size=10.5, color=TEXT_MUTED, line_spacing=1.4)
            cy += Inches(0.35)


def slide_18_close(prs, total):
    slide = new_slide(prs)

    # 배경 장식
    add_circle(slide, Inches(9.5), Inches(-2.5), Inches(6), PRIMARY_SOFT)
    add_circle(slide, Inches(-2), Inches(5), Inches(5), ACCENT_SOFT)

    add_rect(slide, Inches(0.6), Inches(0.55), Inches(0.5), Inches(0.5),
             fill=PRIMARY, corner=0.22)
    add_text(slide, Inches(0.6), Inches(0.55), Inches(0.5), Inches(0.5),
             "B", size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)
    add_text(slide, Inches(1.2), Inches(0.65), Inches(4), Inches(0.3),
             "End of deck", size=12, bold=True, color=TEXT_SUBTLE, font=FONT_EN)

    add_pill(slide, Inches(0.6), Inches(2.3), Inches(2.2), Inches(0.4),
             "LET’S CONNECT", fill=PRIMARY_SOFT, color=PRIMARY, size=11, bold=True)
    add_text(slide, Inches(0.6), Inches(2.85), Inches(12.2), Inches(1.1),
             "응급실이 필요할 때,",
             size=60, bold=True, color=TEXT, line_spacing=1.02)
    add_text(slide, Inches(0.6), Inches(4.0), Inches(12.2), Inches(1.1),
             "바로 연결해 드릴게요.",
             size=60, bold=True, color=PRIMARY, line_spacing=1.02)
    add_text(slide, Inches(0.6), Inches(5.3), Inches(11), Inches(1.0),
             "피드백·협업·파트너십 제안은 언제든 환영합니다.\n"
             "감사합니다.",
             size=16, color=TEXT_MUTED, line_spacing=1.55)

    # 연락처 박스
    add_rect(slide, Inches(0.6), Inches(6.5), Inches(12.2), Inches(0.55),
             fill=SURFACE, line=BORDER, corner=0.3)
    add_text(slide, Inches(0.9), Inches(6.5), Inches(12), Inches(0.55),
             "✉  support@baroer.app     ·     팀 바로응급실     ·     2026",
             size=13, bold=True, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_01_cover,
        slide_02_why,
        slide_03_solution,
        slide_04_features,
        slide_05_team,
        slide_06_process,
        slide_07_stack,
        slide_08_data,
        slide_09_challenges,
        slide_10_compare,
        slide_11_differentiators,
        slide_12_paramedic,
        slide_13_screens,
        slide_14_play_release,
        slide_15_aso,
        slide_16_marketing,
        slide_17_roadmap,
        slide_18_close,
    ]
    total = len(builders)
    for fn in builders:
        fn(prs, total)

    out_dir = os.path.join(os.path.dirname(__file__), "..", "docs")
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.abspath(os.path.join(out_dir, "BaroER_발표자료.pptx"))

    # PowerPoint 로 파일이 열려 있으면 락이 걸려 저장 실패 → timestamp 붙여서 우회
    try:
        prs.save(out_path)
        saved_path = out_path
    except PermissionError:
        import datetime
        ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        saved_path = os.path.abspath(
            os.path.join(out_dir, f"BaroER_발표자료_{ts}.pptx")
        )
        prs.save(saved_path)

    try:
        print(f"[OK] saved: {saved_path}")
        print(f"     slides: {total}  ·  16:9  ·  Malgun Gothic")
    except UnicodeEncodeError:
        pass


if __name__ == "__main__":
    build()
